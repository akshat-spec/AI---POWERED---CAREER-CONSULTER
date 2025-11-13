from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Intelligent Career Guidance System"
    subtitle.text = "A Machine Learning-Powered Career Recommendation Platform\n\nPrepared by: Data Scientist\nDate: 2024"

    # Slide 2: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Introduction'
    tf = body_shape.text_frame
    tf.text = 'Overview of the Intelligent Career Guidance System'

    p = tf.add_paragraph()
    p.text = '• Web-based platform leveraging machine learning for personalized career recommendations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Built with Flask backend and responsive HTML/CSS frontend'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Integrates pre-trained ML model for career prediction based on skill assessments'
    p.level = 1

    # Slide 3: Project Objectives
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Project Objectives'
    tf = body_shape.text_frame
    tf.text = 'Key goals and features of the system'

    p = tf.add_paragraph()
    p.text = '• Create user-friendly web platform for career guidance'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Implement ML-based career prediction algorithm'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Provide seamless authentication and user experience'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Offer additional resources (courses, blogs, contact info)'
    p.level = 1

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'System Architecture'
    tf = body_shape.text_frame
    tf.text = 'Technical components and data flow'

    p = tf.add_paragraph()
    p.text = 'Backend (Flask):'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Routes: /, /signup, /login, /predict, /about, /contact, /blog, /courses'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Authentication: Session-based with in-memory user storage'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ML Integration: Scikit-learn model (careerlast.pkl)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Frontend: HTML templates with Bootstrap styling'
    p.level = 0

    # Slide 5: Machine Learning Model
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Machine Learning Model'
    tf = body_shape.text_frame
    tf.text = 'Details of the career prediction algorithm'

    p = tf.add_paragraph()
    p.text = 'Model Specifications:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Type: Classification model (Random Forest/Ensemble)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Input: 17 skill ratings (1-9 scale)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Output: 17 career categories'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Training Data: dataset9000.csv'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Prediction Process: User ratings → Model → Primary + Secondary careers'
    p.level = 0

    # Slide 6: Key Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Key Features & Functionality'
    tf = body_shape.text_frame
    tf.text = 'Main capabilities of the system'

    p = tf.add_paragraph()
    p.text = 'User Management:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Registration and login system'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Session-based authentication'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Career Assessment:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 17-skill rating form (Database, Programming, AI/ML, etc.)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ML-powered career recommendations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Multiple career suggestions with probability scores'
    p.level = 1

    # Slide 7: User Journey
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'User Journey & Data Flow'
    tf = body_shape.text_frame
    tf.text = 'How users interact with the system'

    p = tf.add_paragraph()
    p.text = '1. Access: Visit http://127.0.0.1:5000'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '2. Registration: Create account with username/password'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '3. Login: Authenticate existing users'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '4. Assessment: Rate skills on 17 categories (1-9 scale)'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '5. Prediction: ML model analyzes input and predicts careers'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '6. Results: View primary and secondary recommendations'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '7. Exploration: Navigate to career details or resources'
    p.level = 0

    # Slide 8: Technologies Used
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Technologies & Tools'
    tf = body_shape.text_frame
    tf.text = 'Tech stack and development environment'

    p = tf.add_paragraph()
    p.text = 'Backend:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Python 3.x, Flask 2.3.3'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• NumPy 1.24.3, Scikit-learn 1.3.0'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Frontend:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• HTML5, CSS3, JavaScript'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Bootstrap, Jinja2 templates'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Development Tools:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• VS Code, Git, Pip'
    p.level = 1

    # Slide 9: Future Enhancements
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Future Enhancements'
    tf = body_shape.text_frame
    tf.text = 'Potential improvements and expansions'

    p = tf.add_paragraph()
    p.text = 'Technical Improvements:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Database integration (PostgreSQL/SQLite)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Advanced ML models (Deep Learning)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• User profile and history tracking'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Feature Additions:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Career roadmap generation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Skill gap analysis'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Job market integration'
    p.level = 1

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Conclusion'
    tf = body_shape.text_frame
    tf.text = 'Summary and impact of the project'

    p = tf.add_paragraph()
    p.text = '• Successfully demonstrates ML application in web development'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Provides valuable career insights through skill assessment'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Showcases end-to-end data science workflow'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Scalable solution for career guidance with room for expansion'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Integrates authentication, ML prediction, and responsive design'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Thank You!'
    p.level = 0

    # Save the presentation
    prs.save('Intelligent_Career_Guidance_System_Presentation.pptx')
    print("Presentation saved as 'Intelligent_Career_Guidance_System_Presentation.pptx'")

    # Convert to .ppt format (older PowerPoint format)
    import os
    from win32com.client import Dispatch

    ppt_app = Dispatch('PowerPoint.Application')
    ppt_app.Visible = False

    pptx_path = os.path.abspath('Intelligent_Career_Guidance_System_Presentation.pptx')
    ppt_path = os.path.abspath('Intelligent_Career_Guidance_System_Presentation.ppt')

    presentation = ppt_app.Presentations.Open(pptx_path)
    presentation.SaveAs(ppt_path, 1)  # 1 = ppSaveAsPresentation (old .ppt format)
    presentation.Close()

    ppt_app.Quit()

    print("Presentation also saved as 'Intelligent_Career_Guidance_System_Presentation.ppt'")

if __name__ == "__main__":
    create_presentation()
